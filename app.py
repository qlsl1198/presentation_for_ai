import json
import os
from io import BytesIO

import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


APP_TITLE = "AI PPT Generator"
DEFAULT_ACCENT = RGBColor(124, 58, 237)
TEXT_COLOR = RGBColor(28, 31, 38)
MUTED_COLOR = RGBColor(92, 98, 112)


def build_fallback_slides(topic, audience, slide_count, tone, language):
    opening = "문제 정의" if language == "한국어" else "Problem"
    market = "시장과 사용자" if language == "한국어" else "Market and Users"
    solution = "핵심 솔루션" if language == "한국어" else "Core Solution"
    ai_value = "AI 활용 방식" if language == "한국어" else "How AI Is Used"
    plan = "구현 계획" if language == "한국어" else "Implementation Plan"
    closing = "기대 효과" if language == "한국어" else "Expected Impact"

    korean = language == "한국어"
    base = [
        {
            "title": topic,
            "subtitle": f"{audience}를 위한 {tone} 발표 자료" if korean else f"A {tone.lower()} deck for {audience}",
            "bullets": [
                "발표 목적과 핵심 메시지를 명확히 전달합니다." if korean else "Clarifies the main purpose and message.",
                "청중이 빠르게 이해할 수 있는 흐름으로 구성합니다." if korean else "Uses a clear flow for fast audience understanding.",
            ],
        },
        {
            "title": opening,
            "subtitle": "왜 이 서비스가 필요한가" if korean else "Why this service matters",
            "bullets": [
                "발표 준비에는 자료 조사, 구조화, 디자인 시간이 많이 필요합니다." if korean else "Presentation work takes time across research, structure, and design.",
                "초보자는 좋은 발표 흐름을 잡기 어렵습니다." if korean else "Beginners often struggle to build a persuasive flow.",
                "반복적인 슬라이드 작성 업무를 자동화할 수 있습니다." if korean else "Repetitive slide drafting can be automated.",
            ],
        },
        {
            "title": market,
            "subtitle": "주요 사용자는 누구인가" if korean else "Who the product serves",
            "bullets": [
                "학생: 과제 발표와 프로젝트 발표를 빠르게 준비합니다." if korean else "Students prepare class and project presentations faster.",
                "창업팀: 투자 제안서와 서비스 소개 자료를 제작합니다." if korean else "Startup teams draft pitch and product decks.",
                "직장인: 회의 보고서와 제안 자료 초안을 생성합니다." if korean else "Professionals create meeting and proposal drafts.",
            ],
        },
        {
            "title": solution,
            "subtitle": "입력 한 번으로 발표 초안 생성" if korean else "Draft a deck from one prompt",
            "bullets": [
                "사용자는 주제, 대상, 톤, 슬라이드 수만 입력합니다." if korean else "Users enter topic, audience, tone, and slide count.",
                "AI가 제목, 소제목, 핵심 bullet을 자동 생성합니다." if korean else "AI generates titles, subtitles, and key bullets.",
                "생성 결과를 PPTX 파일로 다운로드할 수 있습니다." if korean else "The result can be downloaded as a PPTX file.",
            ],
        },
        {
            "title": ai_value,
            "subtitle": "생성형 AI로 발표 구조를 설계" if korean else "Designing the flow with generative AI",
            "bullets": [
                "주제에 맞는 발표 목차를 자동 구성합니다." if korean else "Automatically builds an outline for the topic.",
                "청중 수준에 맞춰 설명 깊이와 어휘를 조절합니다." if korean else "Adjusts depth and wording for the audience.",
                "슬라이드별 메시지를 짧고 명확하게 정리합니다." if korean else "Keeps each slide concise and message-driven.",
            ],
        },
        {
            "title": plan,
            "subtitle": "Streamlit 기반 MVP" if korean else "Streamlit-based MVP",
            "bullets": [
                "입력 화면, AI 생성 로직, PPTX 변환 기능으로 구성합니다." if korean else "Includes input UI, AI generation logic, and PPTX export.",
                "python-pptx를 사용해 발표 파일을 생성합니다." if korean else "Uses python-pptx to create presentation files.",
                "향후 템플릿 선택, 이미지 생성, 편집 기능을 추가할 수 있습니다." if korean else "Can later add templates, image generation, and editing.",
            ],
        },
        {
            "title": closing,
            "subtitle": "발표 준비 시간을 줄이는 생산성 도구" if korean else "A productivity tool for faster presentation work",
            "bullets": [
                "초안 작성 시간을 줄이고 내용 품질을 일정하게 유지합니다." if korean else "Reduces drafting time and keeps content quality consistent.",
                "발표 경험이 적은 사용자도 구조적인 자료를 만들 수 있습니다." if korean else "Helps inexperienced presenters create structured decks.",
                "교육, 창업, 업무 환경에서 확장 가능성이 높습니다." if korean else "Useful across education, startup, and workplace scenarios.",
            ],
        },
    ]

    if slide_count <= len(base):
        return base[:slide_count]

    extra_slides = []
    for index in range(len(base) + 1, slide_count + 1):
        extra_slides.append(
            {
                "title": f"추가 분석 {index}" if korean else f"Additional Insight {index}",
                "subtitle": "발표 내용을 더 구체화하는 보조 슬라이드" if korean else "A supporting slide that adds more detail",
                "bullets": [
                    "핵심 메시지를 하나의 관점으로 정리합니다." if korean else "Frames one focused message.",
                    "사용자 가치와 구현 가능성을 연결합니다." if korean else "Connects user value with implementation feasibility.",
                    "발표 흐름에 맞춰 예시나 근거를 보강합니다." if korean else "Adds examples or evidence to support the flow.",
                ],
            }
        )
    return base + extra_slides


def parse_ai_json(raw_text):
    cleaned = raw_text.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.strip("`")
        cleaned = cleaned.removeprefix("json").strip()
    data = json.loads(cleaned)
    return data["slides"]


def generate_with_openai(topic, audience, slide_count, tone, language):
    from openai import OpenAI

    client = OpenAI()
    prompt = f"""
Create a presentation outline as strict JSON.
Language: {language}
Topic: {topic}
Audience: {audience}
Tone: {tone}
Slide count: {slide_count}

Return only this JSON shape:
{{
  "slides": [
    {{
      "title": "short slide title",
      "subtitle": "short subtitle",
      "bullets": ["bullet 1", "bullet 2", "bullet 3"]
    }}
  ]
}}
"""
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You create concise startup presentation decks."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.7,
    )
    content = response.choices[0].message.content or ""
    return parse_ai_json(content)


def generate_slides(topic, audience, slide_count, tone, language):
    if os.getenv("OPENAI_API_KEY"):
        try:
            return generate_with_openai(topic, audience, slide_count, tone, language), "OpenAI"
        except Exception as error:
            st.warning(f"AI 생성에 실패해서 기본 생성기로 전환했습니다: {error}")
    return build_fallback_slides(topic, audience, slide_count, tone, language), "기본 생성기"


def add_textbox(slide, text, left, top, width, height, font_size, color=TEXT_COLOR, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def style_background(slide, index):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DEFAULT_ACCENT
    accent.line.fill.background()

    marker = slide.shapes.add_shape(1, Inches(12.35), Inches(0.42), Inches(0.52), Inches(0.52))
    marker.fill.solid()
    marker.fill.fore_color.rgb = DEFAULT_ACCENT
    marker.line.fill.background()
    text = marker.text_frame
    text.text = str(index)
    text.paragraphs[0].alignment = PP_ALIGN.CENTER
    text.paragraphs[0].runs[0].font.size = Pt(12)
    text.paragraphs[0].runs[0].font.bold = True
    text.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)


def create_pptx(slides):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for index, slide_data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        style_background(slide, index)

        add_textbox(
            slide,
            slide_data["title"],
            Inches(0.75),
            Inches(0.85),
            Inches(11.1),
            Inches(0.75),
            32,
            bold=True,
        )
        add_textbox(
            slide,
            slide_data.get("subtitle", ""),
            Inches(0.78),
            Inches(1.62),
            Inches(10.8),
            Inches(0.45),
            15,
            color=MUTED_COLOR,
        )

        top = 2.45
        for bullet in slide_data.get("bullets", []):
            dot = slide.shapes.add_shape(9, Inches(0.9), Inches(top + 0.08), Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = DEFAULT_ACCENT
            dot.line.fill.background()
            add_textbox(slide, bullet, Inches(1.18), Inches(top), Inches(10.9), Inches(0.45), 18)
            top += 0.68

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream


def render_slide_preview(slide, index):
    with st.container(border=True):
        st.caption(f"Slide {index}")
        st.subheader(slide["title"])
        if slide.get("subtitle"):
            st.write(slide["subtitle"])
        for bullet in slide.get("bullets", []):
            st.markdown(f"- {bullet}")


def main():
    st.set_page_config(page_title=APP_TITLE, page_icon="📊", layout="wide")
    st.title(APP_TITLE)
    st.write("주제를 입력하면 발표 흐름을 만들고 PowerPoint 파일로 다운로드합니다.")

    with st.sidebar:
        st.header("설정")
        topic = st.text_input("발표 주제", value="AI PPT 생성기")
        audience = st.text_input("청중", value="대학생과 예비 창업자")
        slide_count = st.slider("슬라이드 수", min_value=4, max_value=12, value=7)
        tone = st.selectbox("발표 톤", ["설득력 있는", "전문적인", "친근한", "간결한"])
        language = st.selectbox("언어", ["한국어", "English"])
        generate = st.button("PPT 초안 생성", type="primary", use_container_width=True)

    if "slides" not in st.session_state or generate:
        st.session_state.slides, source = generate_slides(topic, audience, slide_count, tone, language)
        st.session_state.source = source

    slides = st.session_state.slides
    st.info(f"생성 방식: {st.session_state.source}")

    left, right = st.columns([0.62, 0.38])
    with left:
        st.header("생성된 슬라이드")
        for index, slide in enumerate(slides, start=1):
            render_slide_preview(slide, index)

    with right:
        st.header("다운로드")
        pptx_file = create_pptx(slides)
        st.download_button(
            "PowerPoint 다운로드",
            data=pptx_file,
            file_name=f"{topic.replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
        st.caption("다운로드한 파일은 PowerPoint, Keynote, Google Slides에서 열 수 있습니다.")


if __name__ == "__main__":
    main()
